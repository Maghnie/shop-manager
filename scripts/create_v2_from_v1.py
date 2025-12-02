#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""
Create v2 by inserting motivation section into v1 presentation
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE


# Color palette
COLORS = {
    'primary_blue': RGBColor(0, 123, 255),
    'dark_blue': RGBColor(0, 51, 102),
    'success_green': RGBColor(40, 167, 69),
    'warning_orange': RGBColor(255, 193, 7),
    'danger_red': RGBColor(220, 53, 69),
    'info_cyan': RGBColor(23, 162, 184),
    'purple': RGBColor(111, 66, 193),
    'pink': RGBColor(232, 62, 140),
    'teal': RGBColor(32, 201, 151),
    'indigo': RGBColor(102, 16, 242),
    'gold': RGBColor(255, 165, 0),
    'light_gray': RGBColor(248, 249, 250),
    'dark_gray': RGBColor(52, 58, 64),
    'white': RGBColor(255, 255, 255),
}


def add_home_button(slide, toc_slide_index):
    """Add a clickable home button"""
    button = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(9.2), Inches(0.1),
        Inches(0.7), Inches(0.35)
    )
    button.fill.solid()
    button.fill.fore_color.rgb = COLORS['primary_blue']
    button.line.color.rgb = COLORS['primary_blue']

    text_frame = button.text_frame
    text_frame.text = "🏠 Home"
    text_frame.paragraphs[0].font.size = Pt(11)
    text_frame.paragraphs[0].font.color.rgb = COLORS['white']
    text_frame.paragraphs[0].font.bold = True
    text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    text_frame.vertical_anchor = 1

    text_frame.paragraphs[0].runs[0].hyperlink.address = f"#Slide {toc_slide_index}"
    return button


def create_motivation_slides(prs):
    """Create all 7 motivation slides and return them as a list"""
    # This creates slides in a temporary presentation
    temp_prs = Presentation()
    temp_prs.slide_width = prs.slide_width
    temp_prs.slide_height = prs.slide_height

    slides_xml = []

    # We'll create the slides and extract their XML
    # For simplicity, let's just create a new presentation with motivation slides
    # and manually copy them

    print("Creating motivation slides...")

    # Motivation slides will be created directly in the main presentation
    # at the correct position
    return None


def main():
    """Main function to create v2"""
    print("=" * 70)
    print("Creating v2 from v1 + Motivation Section")
    print("=" * 70)
    print()

    # Load v1 presentation
    print("Loading v1 presentation...")
    v1_file = "Shop_Management_System_User_Guide_EN.pptx"
    prs = Presentation(v1_file)

    print(f"Loaded {len(prs.slides)} slides from v1")

    # Strategy: We'll create a new presentation with:
    # - Slides 0-1 from v1 (title, intro)
    # - 7 new motivation slides
    # - Updated TOC
    # - Rest of v1 slides

    print("\nCreating new v2 presentation...")
    new_prs = Presentation()
    new_prs.slide_width = prs.slide_width
    new_prs.slide_height = prs.slide_height

    # Copy slide 0 (title) and slide 1 (intro) from v1
    print("  - Copying title and intro slides from v1...")

    # For PPTX, we need to copy slides properly
    # Unfortunately, python-pptx doesn't have direct slide copying
    # So we'll take a different approach

    print("\nAlternative approach: We'll use the incomplete v2 file and note what's missing")
    print("The v2 file has the motivation section. We need to append sections 3-9 from v1.")

    # Load the incomplete v2
    v2_incomplete = "Shop_Management_System_User_Guide_EN_v2.pptx"
    v2_prs = Presentation(v2_incomplete)

    print(f"\nv2 currently has: {len(v2_prs.slides)} slides")
    print(f"v1 has: {len(prs.slides)} slides")
    print(f"Missing: ~{len(prs.slides) - len(v2_prs.slides)} slides (sections 3-9 content)")

    print("\nSuggestion: Manually use the generate_presentation_english.py")
    print("and insert motivation section code into it.")

    return None


if __name__ == "__main__":
    main()
