#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""
Generate Enhanced English User Guide PowerPoint Presentation for Shop Management System
With colors, emojis, and navigation buttons
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE


# Color palette for the presentation
COLORS = {
    'primary_blue': RGBColor(0, 123, 255),
    'dark_blue': RGBColor(0, 51, 102),
    'success_green': RGBColor(40, 167, 69),
    'warning_orange': RGBColor(255, 193, 7),
    'danger_red': RGBColor(220, 53, 69),
    'info_cyan': RGBColor(23, 162, 184),
    'purple': RGBColor(111, 66, 193),
    'pink': RGBColor(232, 62, 140),
    'teal': RGBColor(32, 201, 151),
    'indigo': RGBColor(102, 16, 242),
    'light_gray': RGBColor(248, 249, 250),
    'dark_gray': RGBColor(52, 58, 64),
    'white': RGBColor(255, 255, 255),
}


def add_home_button(slide, toc_slide_index=2):
    """Add a clickable home button to navigate back to TOC"""
    # Create a rounded rectangle button
    button = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(9.2), Inches(0.1),
        Inches(0.7), Inches(0.35)
    )

    # Style the button
    button.fill.solid()
    button.fill.fore_color.rgb = COLORS['primary_blue']
    button.line.color.rgb = COLORS['primary_blue']
    button.shadow.inherit = False

    # Add text
    text_frame = button.text_frame
    text_frame.text = "🏠 Home"
    text_frame.paragraphs[0].font.size = Pt(11)
    text_frame.paragraphs[0].font.color.rgb = COLORS['white']
    text_frame.paragraphs[0].font.bold = True
    text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
    text_frame.vertical_anchor = 1  # Middle

    # Add hyperlink to TOC slide
    text_frame.paragraphs[0].runs[0].hyperlink.address = f"#Slide {toc_slide_index}"

    return button


def add_title_slide(prs):
    """Create title slide"""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    # Gradient background effect using shapes
    shapes = slide.shapes

    # Top gradient
    bg_top = shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        prs.slide_width, Inches(4)
    )
    bg_top.fill.solid()
    bg_top.fill.fore_color.rgb = COLORS['primary_blue']
    bg_top.line.fill.background()

    # Bottom gradient
    bg_bottom = shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(4),
        prs.slide_width, Inches(3.5)
    )
    bg_bottom.fill.solid()
    bg_bottom.fill.fore_color.rgb = COLORS['dark_blue']
    bg_bottom.line.fill.background()

    # Main title
    title_box = shapes.add_textbox(
        Inches(0.5), Inches(2),
        Inches(9), Inches(1.2)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "🏪 Your Shop Management System"
    p.font.size = Pt(60)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    subtitle_box = shapes.add_textbox(
        Inches(1), Inches(3.5),
        Inches(8), Inches(1.5)
    )
    tf = subtitle_box.text_frame
    p = tf.paragraphs[0]
    p.text = "Complete User Guide\n\n💼 Smart & Integrated System for\nPlastic Consumables Shop Management"
    p.font.size = Pt(26)
    p.font.color.rgb = COLORS['white']
    p.alignment = PP_ALIGN.CENTER
    p.line_spacing = 1.3

    # Decorative elements
    for i in range(3):
        circle = shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(1 + i * 3), Inches(6.2),
            Inches(0.8), Inches(0.8)
        )
        circle.fill.solid()
        colors = [COLORS['success_green'], COLORS['warning_orange'], COLORS['danger_red']]
        circle.fill.fore_color.rgb = colors[i]
        circle.line.fill.background()

    return slide


def add_introduction_slide(prs):
    """Create introduction slide"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    shapes = slide.shapes

    # Background
    bg = shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        prs.slide_width, prs.slide_height
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLORS['light_gray']
    bg.line.fill.background()

    # Title with icon
    title_box = shapes.add_textbox(
        Inches(0.5), Inches(0.5),
        Inches(9), Inches(0.7)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "👋 Welcome to Your Shop Management System"
    p.font.size = Pt(40)
    p.font.bold = True
    p.font.color.rgb = COLORS['dark_blue']
    p.alignment = PP_ALIGN.CENTER

    # Two column content
    # Left column - What is it?
    left_box = shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.5), Inches(1.5),
        Inches(4.3), Inches(5)
    )
    left_box.fill.solid()
    left_box.fill.fore_color.rgb = COLORS['white']
    left_box.line.color.rgb = COLORS['primary_blue']
    left_box.line.width = Pt(2)

    left_text = left_box.text_frame
    left_text.margin_left = Inches(0.2)
    left_text.margin_right = Inches(0.2)
    left_text.margin_top = Inches(0.2)

    content_left = [
        ("🎯 What is this System?", 24, True, COLORS['primary_blue']),
        ("A complete solution for managing your retail business efficiently and professionally", 16, False, COLORS['dark_gray']),
        ("", 12, False, COLORS['dark_gray']),
        ("✨ Key Highlights:", 20, True, COLORS['success_green']),
        ("• Easy to use interface", 16, False, COLORS['dark_gray']),
        ("• Real-time calculations", 16, False, COLORS['dark_gray']),
        ("• Automatic alerts", 16, False, COLORS['dark_gray']),
        ("• Professional invoices", 16, False, COLORS['dark_gray']),
        ("• Detailed analytics", 16, False, COLORS['dark_gray']),
    ]

    for i, (text, size, bold, color) in enumerate(content_left):
        if i == 0:
            p = left_text.paragraphs[0]
        else:
            p = left_text.add_paragraph()
        p.text = text
        p.font.size = Pt(size)
        p.font.bold = bold
        p.font.color.rgb = color
        if text.startswith("•"):
            p.level = 1

    # Right column - Why use it?
    right_box = shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(5.2), Inches(1.5),
        Inches(4.3), Inches(5)
    )
    right_box.fill.solid()
    right_box.fill.fore_color.rgb = COLORS['white']
    right_box.line.color.rgb = COLORS['success_green']
    right_box.line.width = Pt(2)

    right_text = right_box.text_frame
    right_text.margin_left = Inches(0.2)
    right_text.margin_right = Inches(0.2)
    right_text.margin_top = Inches(0.2)

    content_right = [
        ("💡 Why Use This System?", 24, True, COLORS['success_green']),
        ("Transform your business operations with powerful features", 16, False, COLORS['dark_gray']),
        ("", 12, False, COLORS['dark_gray']),
        ("📈 Business Benefits:", 20, True, COLORS['purple']),
        ("• Accurate inventory tracking", 16, False, COLORS['dark_gray']),
        ("• Increased profit visibility", 16, False, COLORS['dark_gray']),
        ("• Reduced stock-outs", 16, False, COLORS['dark_gray']),
        ("• Better customer service", 16, False, COLORS['dark_gray']),
        ("• Data-driven decisions", 16, False, COLORS['dark_gray']),
    ]

    for i, (text, size, bold, color) in enumerate(content_right):
        if i == 0:
            p = right_text.paragraphs[0]
        else:
            p = right_text.add_paragraph()
        p.text = text
        p.font.size = Pt(size)
        p.font.bold = bold
        p.font.color.rgb = color
        if text.startswith("•"):
            p.level = 1

    add_home_button(slide)
    return slide


def add_toc_slide(prs):
    """Create Table of Contents slide"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    shapes = slide.shapes

    # Background gradient
    bg = shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        prs.slide_width, prs.slide_height
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLORS['light_gray']
    bg.line.fill.background()

    # Title
    title_box = shapes.add_textbox(
        Inches(0.5), Inches(0.4),
        Inches(9), Inches(0.8)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "📑 Table of Contents"
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = COLORS['dark_blue']
    p.alignment = PP_ALIGN.CENTER

    # TOC items with colorful boxes
    toc_items = [
        ("1️⃣  Getting Started", COLORS['primary_blue']),
        ("2️⃣  Dashboard Overview", COLORS['info_cyan']),
        ("3️⃣  Product Management", COLORS['success_green']),
        ("4️⃣  Inventory Management", COLORS['warning_orange']),
        ("5️⃣  Sales Management", COLORS['purple']),
        ("6️⃣  Invoice Management", COLORS['pink']),
        ("7️⃣  Customer Management", COLORS['teal']),
        ("8️⃣  Analytics & Reports", COLORS['indigo']),
        ("9️⃣  Tips & Best Practices", COLORS['danger_red']),
    ]

    start_y = 1.5
    box_height = 0.55
    spacing = 0.05

    for i, (item, color) in enumerate(toc_items):
        # Create colored box
        box = shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(1.5), Inches(start_y + i * (box_height + spacing)),
            Inches(7), Inches(box_height)
        )
        box.fill.solid()
        box.fill.fore_color.rgb = color
        box.line.fill.background()
        box.shadow.inherit = False

        # Add text
        tf = box.text_frame
        tf.vertical_anchor = 1  # Middle
        p = tf.paragraphs[0]
        p.text = item
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = COLORS['white']
        p.alignment = PP_ALIGN.LEFT
        tf.margin_left = Inches(0.3)

    # Note at bottom
    note_box = shapes.add_textbox(
        Inches(1), Inches(6.8),
        Inches(8), Inches(0.5)
    )
    tf = note_box.text_frame
    p = tf.paragraphs[0]
    p.text = "💡 Tip: Click the 🏠 Home button on any slide to return to this page"
    p.font.size = Pt(16)
    p.font.italic = True
    p.font.color.rgb = COLORS['dark_gray']
    p.alignment = PP_ALIGN.CENTER

    return slide


def add_section_divider(prs, section_number, section_title, emoji, description="", color=None):
    """Create colorful section divider slide"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    shapes = slide.shapes

    if color is None:
        colors = [
            COLORS['primary_blue'],
            COLORS['info_cyan'],
            COLORS['success_green'],
            COLORS['warning_orange'],
            COLORS['purple'],
            COLORS['pink'],
            COLORS['teal'],
            COLORS['indigo'],
            COLORS['danger_red'],
        ]
        color = colors[(section_number - 1) % len(colors)]

    # Background
    bg = shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        prs.slide_width, prs.slide_height
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.line.fill.background()

    # Decorative circles
    for i in range(5):
        circle = shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(-0.5 + i * 2.5), Inches(6),
            Inches(2), Inches(2)
        )
        circle.fill.solid()
        circle.fill.fore_color.rgb = COLORS['white']
        circle.fill.transparency = 0.9
        circle.line.fill.background()

    # Section emoji (large)
    emoji_box = shapes.add_textbox(
        Inches(0.5), Inches(2),
        Inches(9), Inches(1.5)
    )
    tf = emoji_box.text_frame
    p = tf.paragraphs[0]
    p.text = emoji
    p.font.size = Pt(120)
    p.alignment = PP_ALIGN.CENTER

    # Section number and title
    title_box = shapes.add_textbox(
        Inches(1), Inches(3.8),
        Inches(8), Inches(1.2)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"Section {section_number}\n{section_title}"
    p.font.size = Pt(48)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    p.alignment = PP_ALIGN.CENTER
    p.line_spacing = 1.2

    # Description
    if description:
        desc_box = shapes.add_textbox(
            Inches(2), Inches(5.2),
            Inches(6), Inches(0.8)
        )
        tf = desc_box.text_frame
        p = tf.paragraphs[0]
        p.text = description
        p.font.size = Pt(24)
        p.font.color.rgb = COLORS['white']
        p.alignment = PP_ALIGN.CENTER

    add_home_button(slide)
    return slide


def add_content_slide(prs, title_text, title_emoji, content_items, image_description=None, title_color=None):
    """Create content slide with emoji, bullet points and optional image placeholder"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    shapes = slide.shapes

    if title_color is None:
        title_color = COLORS['dark_blue']

    # Background
    bg = shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        prs.slide_width, prs.slide_height
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLORS['light_gray']
    bg.line.fill.background()

    # Title with emoji
    title_box = shapes.add_textbox(
        Inches(0.5), Inches(0.3),
        Inches(9), Inches(0.7)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"{title_emoji} {title_text}"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = title_color

    # Content area
    if image_description:
        # Content with image placeholder
        content_box = shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.5), Inches(1.2),
            Inches(5.5), Inches(5.8)
        )
        content_box.fill.solid()
        content_box.fill.fore_color.rgb = COLORS['white']
        content_box.line.color.rgb = title_color
        content_box.line.width = Pt(2)

        # Image placeholder
        img_box = shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(6.3), Inches(1.2),
            Inches(3.2), Inches(5.8)
        )
        img_box.fill.solid()
        img_box.fill.fore_color.rgb = COLORS['white']
        img_box.line.color.rgb = COLORS['info_cyan']
        img_box.line.width = Pt(2)
        img_box.line.dash_style = 2  # Dashed line

        # Image placeholder icon
        img_icon = shapes.add_textbox(
            Inches(6.3), Inches(2.5),
            Inches(3.2), Inches(1)
        )
        tf_img = img_icon.text_frame
        p_img = tf_img.paragraphs[0]
        p_img.text = "📸"
        p_img.font.size = Pt(80)
        p_img.alignment = PP_ALIGN.CENTER

        # Image description
        img_desc = shapes.add_textbox(
            Inches(6.4), Inches(4),
            Inches(3), Inches(2.5)
        )
        tf_desc = img_desc.text_frame
        tf_desc.word_wrap = True
        p_desc = tf_desc.paragraphs[0]
        p_desc.text = f"[Screenshot]\n\n{image_description}"
        p_desc.font.size = Pt(13)
        p_desc.font.italic = True
        p_desc.alignment = PP_ALIGN.CENTER
        p_desc.font.color.rgb = COLORS['dark_gray']

        content_width = 5.5
    else:
        # Full width content
        content_box = shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(0.5), Inches(1.2),
            Inches(9), Inches(5.8)
        )
        content_box.fill.solid()
        content_box.fill.fore_color.rgb = COLORS['white']
        content_box.line.color.rgb = title_color
        content_box.line.width = Pt(2)
        content_width = 9

    # Add content text
    text_box = shapes.add_textbox(
        Inches(0.7), Inches(1.4),
        Inches(content_width - 0.4), Inches(5.4)
    )
    tf = text_box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.2)
    tf.margin_right = Inches(0.2)

    for i, item in enumerate(content_items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        if isinstance(item, dict):
            p.text = item['text']
            p.level = item.get('level', 0)
            p.font.size = Pt(item.get('size', 18))
            p.font.bold = item.get('bold', False)
            color = item.get('color', COLORS['dark_gray'])
            p.font.color.rgb = color
        else:
            p.text = item
            p.font.size = Pt(17)
            p.font.color.rgb = COLORS['dark_gray']
            if item.startswith("•") or item.startswith("✓") or item.startswith("❌"):
                p.level = 1

    add_home_button(slide)
    return slide


def add_two_column_slide(prs, title_text, title_emoji, left_items, right_items, left_title="", right_title="", left_emoji="", right_emoji=""):
    """Create slide with two columns of content"""
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    shapes = slide.shapes

    # Background
    bg = shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        prs.slide_width, prs.slide_height
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLORS['light_gray']
    bg.line.fill.background()

    # Title with emoji
    title_box = shapes.add_textbox(
        Inches(0.5), Inches(0.3),
        Inches(9), Inches(0.7)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"{title_emoji} {title_text}"
    p.font.size = Pt(36)
    p.font.bold = True
    p.font.color.rgb = COLORS['dark_blue']

    # Left column
    left_box = shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.5), Inches(1.3),
        Inches(4.6), Inches(5.7)
    )
    left_box.fill.solid()
    left_box.fill.fore_color.rgb = COLORS['white']
    left_box.line.color.rgb = COLORS['success_green']
    left_box.line.width = Pt(3)

    left_text = shapes.add_textbox(
        Inches(0.7), Inches(1.5),
        Inches(4.2), Inches(5.3)
    )
    tf_left = left_text.text_frame
    tf_left.word_wrap = True

    if left_title:
        p = tf_left.paragraphs[0]
        p.text = f"{left_emoji} {left_title}"
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = COLORS['success_green']
        p.space_after = Pt(10)

    for item in left_items:
        p = tf_left.add_paragraph()
        p.text = item
        p.font.size = Pt(16)
        p.font.color.rgb = COLORS['dark_gray']
        if item.startswith("•"):
            p.level = 1

    # Right column
    right_box = shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(5.4), Inches(1.3),
        Inches(4.6), Inches(5.7)
    )
    right_box.fill.solid()
    right_box.fill.fore_color.rgb = COLORS['white']
    right_box.line.color.rgb = COLORS['warning_orange']
    right_box.line.width = Pt(3)

    right_text = shapes.add_textbox(
        Inches(5.6), Inches(1.5),
        Inches(4.2), Inches(5.3)
    )
    tf_right = right_text.text_frame
    tf_right.word_wrap = True

    if right_title:
        p = tf_right.paragraphs[0]
        p.text = f"{right_emoji} {right_title}"
        p.font.size = Pt(24)
        p.font.bold = True
        p.font.color.rgb = COLORS['warning_orange']
        p.space_after = Pt(10)

    for item in right_items:
        p = tf_right.add_paragraph()
        p.text = item
        p.font.size = Pt(16)
        p.font.color.rgb = COLORS['dark_gray']
        if item.startswith("•"):
            p.level = 1

    add_home_button(slide)
    return slide


def generate_presentation():
    """Generate the complete enhanced English presentation"""
    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    print("Creating enhanced English presentation with colors and navigation...")
    print()

    # 1. Title Slide
    print("  - Title slide")
    add_title_slide(prs)

    # 2. Introduction
    print("  - Introduction")
    add_introduction_slide(prs)

    # 3. Table of Contents (slide index 3, but 0-based is 2)
    print("  - Table of Contents")
    add_toc_slide(prs)

    # SECTION 1: Getting Started
    print("  - Section 1: Getting Started")
    add_section_divider(prs, 1, "Getting Started", "🚀", "Your quick guide to begin using the system")

    add_content_slide(prs, "Accessing the System", "🔐", [
        "How to log into your system:",
        "",
        "1️⃣ Open your web browser (Chrome, Firefox, Edge, or Safari)",
        "",
        "2️⃣ Navigate to your system's URL",
        "",
        "3️⃣ Enter your username and password",
        "",
        "4️⃣ Click the 'Login' button",
        "",
        "🔒 Security Tips:",
        "• Keep your password secure and confidential",
        "• Never share your login credentials",
        "• Always log out when finished",
        "• Use a strong, unique password"
    ], "Screenshot showing the login screen with username and password fields, and login button", COLORS['primary_blue'])

    add_content_slide(prs, "First Steps After Login", "👣", [
        "What you'll see after logging in:",
        "",
        "📊 Main Dashboard:",
        "• Quick overview of today's sales",
        "• Total revenue and profit metrics",
        "• Low stock alerts and warnings",
        "• Key performance indicators",
        "",
        "🧭 Main Navigation Menu:",
        "• Home (Dashboard)",
        "• Products",
        "• Inventory",
        "• Sales",
        "• Analytics & Reports",
        "",
        "Take a moment to familiarize yourself with the layout!"
    ], None, COLORS['primary_blue'])

    add_content_slide(prs, "System Overview", "🗺️", [
        "Understanding the big picture:",
        "",
        "The system is divided into logical sections:",
        "",
        "📦 Product Management - Your catalog",
        "📊 Inventory Tracking - Stock levels",
        "💰 Sales Processing - Transactions",
        "📄 Invoicing - Professional receipts",
        "👥 Customer Management - Client database",
        "📈 Analytics - Business insights",
        "",
        "💡 Pro Tip: Start by setting up your products, then add inventory, and you're ready to make sales!"
    ], None, COLORS['primary_blue'])

    # SECTION 2: Dashboard
    print("  - Section 2: Dashboard")
    add_section_divider(prs, 2, "Dashboard Overview", "📊", "Your business at a glance")

    add_content_slide(prs, "Understanding Your Dashboard", "🎯", [
        "Your command center for daily operations:",
        "",
        "📈 Key Metrics (At a Glance):",
        "• Today's sales count",
        "• Today's total revenue",
        "• Today's net profit & profit margin",
        "• This month's sales & revenue",
        "",
        "🎨 Color-Coded Cards:",
        "• Blue cards - Sales information",
        "• Green cards - Profit data",
        "• Orange cards - Monthly stats",
        "• Red cards - Alerts & warnings",
        "",
        "⚡ The dashboard updates in real-time as you make sales!"
    ], "Screenshot of the dashboard showing all KPI cards with sales, revenue, and profit metrics", COLORS['info_cyan'])

    add_content_slide(prs, "Quick Action Buttons", "⚡", [
        "Speed up your workflow with shortcuts:",
        "",
        "🛒 New Sale Button:",
        "• Jump directly to creating a sale",
        "• Perfect for busy times",
        "• One-click access to sales form",
        "",
        "📦 Manage Products:",
        "• View and edit product catalog",
        "• Add new products quickly",
        "",
        "📄 View Invoices:",
        "• Access all generated invoices",
        "• Print customer receipts",
        "",
        "⚠️ Low Stock Alerts:",
        "• See products running low",
        "• Prevent stock-outs",
        "• Plan re-orders efficiently"
    ], "Screenshot showing the quick action button section with colorful buttons", COLORS['info_cyan'])

    # SECTION 3: Product Management
    print("  - Section 3: Product Management")
    add_section_divider(prs, 3, "Product Management", "📦", "Organize your entire product catalog")

    add_content_slide(prs, "Why Product Management Matters", "💼", [
        "Well-managed products = Higher profits",
        "",
        "✨ Benefits of organized products:",
        "• Accurate pricing and profit tracking",
        "• Easy searching and finding items",
        "• Better inventory control",
        "• Professional invoices",
        "• Data-driven product decisions",
        "",
        "🎯 What you can do:",
        "• Add new products with full details",
        "• Edit existing product information",
        "• Archive discontinued items (non-destructive)",
        "• Search and filter your catalog",
        "• Organize by type, brand, and material",
        "",
        "💡 Tip: Spend time organizing your products well upfront - it pays off later!"
    ], None, COLORS['success_green'])

    add_content_slide(prs, "Adding a New Product - Step by Step", "➕", [
        "Creating products is quick and easy:",
        "",
        "1️⃣ Click 'Add New Product' button",
        "",
        "2️⃣ Fill in Required Information:",
        "• Product Type (cups, plates, bags, etc.)",
        "• Product name (English and Arabic)",
        "• Cost Price (what you paid)",
        "• Selling Price (what you charge)",
        "",
        "3️⃣ Optional Details (Recommended):",
        "• Brand name",
        "• Material (plastic type, paper, etc.)",
        "• Size and weight",
        "• Tags for easy searching",
        "",
        "4️⃣ Save - and you're done! ✅"
    ], "Screenshot of the add product form with all fields visible and highlighted", COLORS['success_green'])

    add_content_slide(prs, "Pricing and Profit Calculations", "💰", [
        "The system does the math for you:",
        "",
        "📊 Required Pricing Information:",
        "• Cost Price: What you paid to supplier",
        "• Selling Price: What you charge customers",
        "",
        "🧮 Automatic Calculations:",
        "• Profit per unit = Selling Price - Cost Price",
        "• Profit Margin % = (Profit ÷ Cost) × 100",
        "",
        "✅ Built-in Validation:",
        "• System prevents selling below cost",
        "• Warnings for unusual prices",
        "• Ensures profitability",
        "",
        "💡 Example: Cost $10, Sell $15 → $5 profit (50% margin)"
    ], None, COLORS['success_green'])

    add_content_slide(prs, "Organizing Your Products", "🗂️", [
        "Multiple ways to categorize:",
        "",
        "🏷️ By Product Type:",
        "• Plastic cups",
        "• Plates and utensils",
        "• Trash bags",
        "• Food packaging",
        "• Container lids",
        "",
        "🏢 By Brand:",
        "• Organize by manufacturer",
        "• Track brand performance",
        "",
        "🧪 By Material:",
        "• PP (Polypropylene)",
        "• PS (Polystyrene)",
        "• PET, Paper, Aluminum",
        "",
        "🔖 Custom Tags: Add your own keywords for quick searches!"
    ], "Screenshot showing product list with filters for type, brand, and material", COLORS['success_green'])

    add_content_slide(prs, "Search and Filter Tools", "🔍", [
        "Find any product in seconds:",
        "",
        "🔎 Search Options:",
        "• Search by product name",
        "• Filter by type",
        "• Filter by brand",
        "• Filter by material",
        "• Search using tags",
        "",
        "👁️ View Options:",
        "• Show active products only",
        "• Show archived products",
        "• Show all products",
        "",
        "📑 Sorting:",
        "• Sort by name (A-Z)",
        "• Sort by price (low to high)",
        "• Sort by profit margin",
        "",
        "💡 Pro tip: Use tags like 'bestseller' or 'seasonal' for quick filtering!"
    ], "Screenshot of search bar and filter dropdowns in action", COLORS['success_green'])

    # SECTION 4: Inventory Management
    print("  - Section 4: Inventory Management")
    add_section_divider(prs, 4, "Inventory Management", "📊", "Smart tracking of your stock levels")

    add_content_slide(prs, "The Importance of Inventory Control", "⚠️", [
        "Your inventory is your livelihood:",
        "",
        "❌ Problems Without Good Inventory Management:",
        "• Stock-outs during peak times",
        "• Lost sales and disappointed customers",
        "• Overstocking slow-moving items",
        "• Cash tied up in excess inventory",
        "• Expired or damaged stock",
        "",
        "✅ With Our System:",
        "• Real-time stock tracking",
        "• Automatic low-stock alerts",
        "• Instant updates after each sale",
        "• Clear visibility of all products",
        "• Prevention of overselling",
        "",
        "💰 Result: Better cash flow and happier customers!"
    ], None, COLORS['warning_orange'])

    add_content_slide(prs, "Understanding Stock Alerts", "🚨", [
        "Two types of alerts to watch:",
        "",
        "🟡 LOW STOCK WARNING:",
        "• Quantity has reached your minimum threshold",
        "• Time to reorder from suppliers",
        "• Product is still available for sale",
        "• Plan ahead to avoid stock-out",
        "",
        "🔴 OUT OF STOCK CRITICAL:",
        "• Quantity is zero",
        "• Cannot sell this product",
        "• Urgent reorder needed",
        "• Customer orders may be affected",
        "",
        "⚙️ You can set custom minimum levels for each product!",
        "",
        "💡 Best Practice: Set minimums to give you 3-5 days to restock"
    ], "Screenshot of inventory page showing products with yellow and red alert indicators", COLORS['warning_orange'])

    add_content_slide(prs, "Automatic Inventory Updates", "🔄", [
        "Let the system handle the counting:",
        "",
        "✅ When You Make a Sale:",
        "• Quantities automatically deducted",
        "• Inventory status updated instantly",
        "• Alerts triggered if needed",
        "• Accurate stock levels maintained",
        "",
        "↩️ When You Cancel a Sale:",
        "• Quantities automatically restored",
        "• Inventory adjusted back",
        "• Status recalculated",
        "",
        "🛡️ Safety Features:",
        "• Cannot sell more than available",
        "• Warning before archiving products with stock",
        "• Validation on all quantity changes",
        "",
        "🎯 You focus on selling, we handle the counting!"
    ], None, COLORS['warning_orange'])

    add_content_slide(prs, "Best Practices for Inventory", "✨", [
        "Pro tips for inventory success:",
        "",
        "📅 Daily Routine:",
        "• Check inventory dashboard each morning",
        "• Review low-stock alerts",
        "• Place orders before critical levels",
        "",
        "📊 Weekly Analysis:",
        "• Identify fast-moving products",
        "• Spot slow-moving inventory",
        "• Adjust minimum stock levels",
        "",
        "🎯 Smart Strategies:",
        "• Keep more stock of bestsellers",
        "• Reduce slow-movers",
        "• Use analytics to predict demand",
        "• Plan for seasonal variations",
        "",
        "💡 The system provides data - you make smart decisions!"
    ], None, COLORS['warning_orange'])

    # SECTION 5: Sales Management
    print("  - Section 5: Sales Management")
    add_section_divider(prs, 5, "Sales Management", "💰", "Record and track every transaction")

    add_content_slide(prs, "Why Record Every Sale?", "📝", [
        "Every recorded sale = Valuable data:",
        "",
        "💡 Benefits:",
        "• Accurate revenue and profit tracking",
        "• Understand product performance",
        "• Professional invoices for customers",
        "• Detailed reports for better decisions",
        "• Track customer purchase history",
        "• Tax and accounting ready",
        "",
        "⚡ Features:",
        "• Fast, simple interface",
        "• Automatic calculations",
        "• Multiple payment methods",
        "• Works with or without customer info",
        "• Real-time profit visibility",
        "",
        "🎯 Make recording sales a habit - your future self will thank you!"
    ], None, COLORS['purple'])

    add_content_slide(prs, "Creating a New Sale - Complete Guide", "🛒", [
        "Selling is easy with our 3-step process:",
        "",
        "👤 STEP 1: Customer Information",
        "• Select existing customer from list, OR",
        "• Enter name and phone directly, OR",
        "• Leave blank for walk-in customers",
        "",
        "📦 STEP 2: Add Products",
        "• Select product from dropdown",
        "• Enter quantity needed",
        "• System calculates total automatically",
        "• Add multiple products as needed",
        "",
        "💳 STEP 3: Finalize & Save",
        "• Review totals and profit",
        "• Choose payment method",
        "• Apply discount if needed (optional)",
        "• Add tax if required (optional)",
        "• Click Save - Done! ✅"
    ], "Screenshot of new sale form showing customer section, product selector, and payment options", COLORS['purple'])

    add_content_slide(prs, "Discounts and Taxes", "💸", [
        "Flexible pricing options:",
        "",
        "🎁 Discount Options:",
        "• Fixed amount discount (e.g., $10 off)",
        "• Percentage discount (e.g., 5% off)",
        "• Applied to subtotal before tax",
        "",
        "🧾 Tax Handling:",
        "• Added as percentage",
        "• Calculated after discount",
        "• Automatically included in total",
        "",
        "🧮 Calculation Formula:",
        "Subtotal = Sum of all items",
        "After Discount = Subtotal - Discount",
        "Tax Amount = After Discount × Tax%",
        "Final Total = After Discount + Tax",
        "",
        "💰 Profit calculated correctly with discounts and taxes!"
    ], None, COLORS['purple'])

    add_content_slide(prs, "Payment Methods", "💳", [
        "Accept payments your way:",
        "",
        "💵 CASH",
        "• Most common method",
        "• Perfect for walk-in customers",
        "• Immediate payment",
        "",
        "💳 CARD (Credit/Debit)",
        "• Card payments",
        "• Secure transactions",
        "",
        "🏦 BANK TRANSFER",
        "• For larger amounts",
        "• Business customers",
        "",
        "📋 CREDIT (Pay Later)",
        "• For trusted customers",
        "• Payment due later",
        "• Track outstanding amounts",
        "",
        "📊 Track all payment types in your reports!"
    ], "Screenshot showing payment method selector with all four options", COLORS['purple'])

    add_content_slide(prs, "Managing Recorded Sales", "📂", [
        "Full control over your sales:",
        "",
        "👀 View & Search:",
        "• See all sales in one list",
        "• Search by sale number",
        "• Search by customer name or phone",
        "• Filter by date range",
        "• Filter by payment method",
        "• Filter by status (completed, pending, cancelled)",
        "",
        "🔧 Actions You Can Take:",
        "• View full details of any sale",
        "• Edit pending sales",
        "• Cancel sales (inventory restored)",
        "• Print invoice",
        "• Email receipt to customer",
        "",
        "🔢 Each sale gets a unique number: S20251118001"
    ], "Screenshot of sales list with search and filter options visible", COLORS['purple'])

    # SECTION 6: Invoice Management
    print("  - Section 6: Invoice Management")
    add_section_divider(prs, 6, "Invoice Management", "📄", "Professional receipts for every sale")

    add_content_slide(prs, "Automatic Invoice Generation", "⚡", [
        "Professional invoices without extra work:",
        "",
        "🎯 Automatic Creation:",
        "• Invoice created when sale is completed",
        "• Unique invoice number assigned",
        "• Timestamp automatically recorded",
        "• Linked to the sale permanently",
        "",
        "📋 What's Included:",
        "• Your company information (name, address, phone, email)",
        "• Customer details",
        "• Itemized product list with quantities",
        "• Unit prices and line totals",
        "• Discounts applied",
        "• Tax calculations",
        "• Final total amount",
        "",
        "🎨 Professional format ready for printing!"
    ], "Screenshot of a sample invoice showing all sections clearly", COLORS['pink'])

    add_content_slide(prs, "Printing and Managing Invoices", "🖨️", [
        "Easy invoice handling:",
        "",
        "🖨️ Printing:",
        "• Print directly from system",
        "• Professional layout",
        "• Print status tracked",
        "• Timestamp when printed",
        "",
        "🔍 Finding Invoices:",
        "• View all invoices in list",
        "• Search by invoice number",
        "• Filter by date",
        "• See print status",
        "• Link to original sale",
        "",
        "⚙️ Customization:",
        "• Add your company logo",
        "• Set company details",
        "• Optional due date",
        "• Custom footer notes",
        "",
        "💡 Tip: Print invoices for all customers - it looks professional!"
    ], None, COLORS['pink'])

    # SECTION 7: Customer Management
    print("  - Section 7: Customer Management")
    add_section_divider(prs, 7, "Customer Management", "👥", "Build relationships with your customers")

    add_content_slide(prs, "Why Maintain Customer Records?", "📇", [
        "Know your customers, grow your business:",
        "",
        "✨ Benefits:",
        "• Faster checkout for repeat customers",
        "• Track purchase history",
        "• Identify your best customers",
        "• Build loyalty",
        "• Credit management for trusted customers",
        "• Marketing opportunities",
        "",
        "🎯 Two Approaches:",
        "• Full customer profiles for regulars",
        "• Quick walk-in sales for one-time customers",
        "",
        "💼 For Plastic Consumables:",
        "• Business customers buy regularly",
        "• Build relationships for repeat orders",
        "• Offer credit to reliable clients",
        "",
        "💡 The more you know, the better you serve!"
    ], None, COLORS['teal'])

    add_content_slide(prs, "Adding a New Customer", "➕", [
        "Quick customer registration:",
        "",
        "📝 Basic Information:",
        "• Full name (Arabic and English)",
        "• Phone number",
        "• Email address",
        "",
        "📍 Optional Details:",
        "• Physical address",
        "• Birth date (for special offers)",
        "• Gender",
        "• Custom notes",
        "",
        "📊 Automatic Statistics:",
        "The system tracks automatically:",
        "• Total amount spent",
        "• Number of purchases",
        "• Date of last purchase",
        "• Average order value",
        "",
        "⚡ Add during sale or beforehand!"
    ], "Screenshot of add customer form with all fields", COLORS['teal'])

    add_two_column_slide(prs, "Registered vs Walk-In Customers", "👥",
        [  # Left column - Registered
            "Complete customer information saved",
            "• Full contact details",
            "• Purchase history tracked",
            "• Accurate invoicing",
            "• Credit sales possible",
            "• Marketing and communication",
            "• Loyalty programs",
            "",
            "🎯 Best For:",
            "• Regular customers",
            "• Business clients",
            "• Credit purchases",
            "• Large orders",
            "• Relationship building"
        ],
        [  # Right column - Walk-in
            "Quick sales without registration",
            "• No details required",
            "• Fast checkout",
            "• Perfect for small sales",
            "• Cash transactions",
            "• High volume periods",
            "",
            "🎯 Best For:",
            "• One-time purchases",
            "• Small transactions",
            "• Busy rush hours",
            "• Cash-only sales",
            "• Anonymous customers"
        ],
        "Registered Customers", "Walk-In Customers", "✅", "🏃"
    )

    # SECTION 8: Analytics & Reports
    print("  - Section 8: Analytics & Reports")
    add_section_divider(prs, 8, "Analytics & Reports", "📈", "Turn data into profits")

    add_content_slide(prs, "The Power of Business Analytics", "💪", [
        "Data-driven decisions = Higher profits:",
        "",
        "🎯 What Analytics Provide:",
        "• Understand your true business performance",
        "• Identify most profitable products",
        "• Spot trends and patterns",
        "• Discover slow-moving inventory",
        "• Plan for peak and slow periods",
        "• Make informed purchase decisions",
        "",
        "📊 Types of Reports:",
        "• Time Series: Track performance over time",
        "• Breakeven Analysis: Product profitability",
        "• Financial Summaries: Overall business health",
        "",
        "💡 Spend 15 minutes weekly reviewing reports - it pays off!",
        "",
        "🚀 From guessing to knowing!"
    ], None, COLORS['indigo'])

    add_content_slide(prs, "Time Series Reports", "📉", [
        "See your performance over time:",
        "",
        "🕐 Available Time Ranges:",
        "• Hourly (up to 7 days)",
        "• Daily (up to 1 year)",
        "• Weekly (up to 2 years)",
        "• Monthly (up to 5 years)",
        "• Yearly (up to 10 years)",
        "",
        "📊 Metrics Tracked:",
        "• Revenue (total sales)",
        "• Costs (COGS)",
        "• Profit (revenue - costs)",
        "• Number of sales",
        "• Profit margin percentage",
        "• Average sale value",
        "",
        "📈 Visual Charts: Easy-to-understand graphs powered by Chart.js"
    ], "Screenshot of time series chart showing revenue and profit trends over time", COLORS['indigo'])

    add_content_slide(prs, "Breakeven Analysis", "⚖️", [
        "Which products deserve your focus?",
        "",
        "🎯 Product-Level Analysis:",
        "• Unit price and cost",
        "• Profit per unit sold",
        "• Profit margin percentage",
        "• Breakeven point (units needed to cover fixed costs)",
        "",
        "📊 Actual Performance:",
        "• Quantity sold to date",
        "• Revenue generated",
        "• Total cost incurred",
        "• Total profit earned",
        "• Performance score vs breakeven",
        "",
        "🏆 Performance Ratings:",
        "🟢 Excellent • 🔵 Good • 🟡 Moderate • 🟠 Profitable • 🔴 Poor",
        "",
        "💡 Focus on excellent performers, improve or drop poor ones!"
    ], "Screenshot of breakeven analysis table showing products with performance ratings", COLORS['indigo'])

    add_content_slide(prs, "Understanding Performance Ratings", "⭐", [
        "What do the ratings mean?",
        "",
        "🟢 EXCELLENT Performance:",
        "• Sales far exceed breakeven point",
        "• Keep promoting these products",
        "• Consider increasing stock",
        "",
        "🔵 GOOD Performance:",
        "• Solid sales and profitability",
        "• Maintain current strategy",
        "",
        "🟡 MODERATE Performance:",
        "• Acceptable but could improve",
        "• Try promotions or better positioning",
        "",
        "🟠 PROFITABLE:",
        "• Making money but below expectations",
        "• Analyze why underperforming",
        "",
        "🔴 POOR Performance:",
        "• Below breakeven or minimal profit",
        "• Consider discontinuing or heavy promotion"
    ], None, COLORS['indigo'])

    add_content_slide(prs, "Exporting Reports", "💾", [
        "Take your data anywhere:",
        "",
        "📁 Export Formats:",
        "• CSV: Open in Excel, Google Sheets, etc.",
        "• XLSX: Native Excel format",
        "• Full data export with all details",
        "",
        "🌐 Arabic Support:",
        "✅ Exported files fully support Arabic text",
        "✅ No encoding issues",
        "✅ Ready for external analysis",
        "",
        "🎯 Use Cases:",
        "• Share with your accountant",
        "• External analysis and modeling",
        "• Backup your data",
        "• Create presentations",
        "• Tax and compliance reporting",
        "",
        "💡 Export reports monthly for your records!"
    ], "Screenshot showing export buttons for CSV and XLSX formats", COLORS['indigo'])

    add_content_slide(prs, "Using Reports to Boost Profits", "🚀", [
        "Turn insights into action:",
        "",
        "📈 Strategy 1: Focus on Winners",
        "• Identify products with excellent ratings",
        "• Increase stock levels",
        "• Give prominent display",
        "• Consider raising prices slightly",
        "",
        "🔄 Strategy 2: Fix Underperformers",
        "• Run promotions on moderate performers",
        "• Bundle slow movers with bestsellers",
        "• Adjust pricing",
        "",
        "❌ Strategy 3: Cut Losers",
        "• Reduce orders for poor performers",
        "• Clear out with sales",
        "• Replace with better alternatives",
        "",
        "📅 Strategy 4: Plan Seasonally",
        "• Review yearly patterns",
        "• Stock up before peak seasons"
    ], None, COLORS['indigo'])

    # SECTION 9: Tips & Best Practices
    print("  - Section 9: Tips & Best Practices")
    add_section_divider(prs, 9, "Tips & Best Practices", "💡", "Master the system like a pro")

    add_content_slide(prs, "Daily Routine for Success", "📅", [
        "Build habits that build profits:",
        "",
        "🌅 Morning (Start of Day):",
        "• Log in and review dashboard",
        "• Check inventory alerts",
        "• Review any pending orders",
        "• Set daily sales goals",
        "",
        "☀️ During the Day:",
        "• Record every sale immediately",
        "• Print invoices for customers",
        "• Monitor product movement",
        "• Note customer feedback",
        "",
        "🌙 Evening (End of Day):",
        "• Review daily sales summary",
        "• Check profit achieved",
        "• Plan restocking needs",
        "• Celebrate wins! 🎉"
    ], None, COLORS['danger_red'])

    add_content_slide(prs, "Weekly and Monthly Reviews", "📊", [
        "Bigger picture analysis:",
        "",
        "📅 Weekly Tasks:",
        "• Review sales report for the week",
        "• Identify products needing restock",
        "• Review customer credit accounts",
        "• Analyze which days were busiest",
        "• Adjust staffing if needed",
        "",
        "📆 Monthly Tasks:",
        "• Export financial reports",
        "• Run breakeven analysis",
        "• Review product performance ratings",
        "• Update product strategy",
        "• Plan next month's inventory",
        "• Review with accountant if applicable",
        "",
        "🎯 Consistency is key to long-term success!"
    ], None, COLORS['danger_red'])

    add_content_slide(prs, "Security and Data Protection", "🔒", [
        "Protect your business data:",
        "",
        "🛡️ Account Security:",
        "• Never share your login credentials",
        "• Use a strong, unique password",
        "• Change password regularly",
        "• Log out after every session",
        "• Don't save passwords in shared computers",
        "",
        "💾 Data Protection:",
        "• Export reports regularly as backups",
        "• Keep copies of important invoices",
        "• Review sales history periodically",
        "• Don't delete data unnecessarily",
        "",
        "❌ When Issues Arise:",
        "• Use archive instead of delete",
        "• Contact support before making major changes",
        "• Keep transaction history intact"
    ], None, COLORS['danger_red'])

    add_content_slide(prs, "Common Mistakes to Avoid", "⚠️", [
        "Learn from others' mistakes:",
        "",
        "❌ DON'T: Delay recording sales",
        "✅ DO: Enter every sale immediately",
        "",
        "❌ DON'T: Ignore inventory alerts",
        "✅ DO: Check alerts daily and reorder promptly",
        "",
        "❌ DON'T: Neglect reports and analytics",
        "✅ DO: Review reports weekly for insights",
        "",
        "❌ DON'T: Enter incorrect prices",
        "✅ DO: Double-check pricing before saving",
        "",
        "❌ DON'T: Skip printing invoices",
        "✅ DO: Give every customer a professional receipt",
        "",
        "💡 Small daily habits lead to big success!"
    ], None, COLORS['danger_red'])

    add_content_slide(prs, "Pro Tips for Power Users", "⚡", [
        "Level up your system mastery:",
        "",
        "🎯 Efficiency Tips:",
        "• Use keyboard shortcuts when available",
        "• Bookmark frequently used pages",
        "• Set appropriate minimum stock levels",
        "• Use product tags extensively",
        "• Create customer profiles for regulars",
        "",
        "📊 Analysis Tips:",
        "• Compare week-over-week performance",
        "• Track profit margins by product category",
        "• Monitor payment method trends",
        "• Identify your top 20% products (80/20 rule)",
        "",
        "💰 Profit Maximization:",
        "• Bundle slow movers with bestsellers",
        "• Offer volume discounts strategically",
        "• Use reports to negotiate with suppliers"
    ], None, COLORS['danger_red'])

    # Closing Slide
    print("  - Closing slide")
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    shapes = slide.shapes

    # Background
    bg = shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(0), Inches(0),
        prs.slide_width, prs.slide_height
    )
    bg.fill.solid()
    bg.fill.fore_color.rgb = COLORS['dark_blue']
    bg.line.fill.background()

    # Decorative circles
    for i in range(6):
        circle = shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(i * 2 - 0.5), Inches(6),
            Inches(1.5), Inches(1.5)
        )
        circle.fill.solid()
        colors_list = [COLORS['primary_blue'], COLORS['success_green'], COLORS['warning_orange'],
                      COLORS['danger_red'], COLORS['purple'], COLORS['teal']]
        circle.fill.fore_color.rgb = colors_list[i % len(colors_list)]
        circle.fill.transparency = 0.7
        circle.line.fill.background()

    # Main message
    title_box = shapes.add_textbox(
        Inches(1), Inches(2),
        Inches(8), Inches(1.5)
    )
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.text = "🎉 Thank You!"
    p.font.size = Pt(72)
    p.font.bold = True
    p.font.color.rgb = COLORS['white']
    p.alignment = PP_ALIGN.CENTER

    # Subtitle
    subtitle_box = shapes.add_textbox(
        Inches(1), Inches(3.8),
        Inches(8), Inches(1.8)
    )
    tf = subtitle_box.text_frame
    p = tf.paragraphs[0]
    p.text = "You're Now Ready to Master\nYour Shop Management System!\n\n📞 Need help? Contact our support team\n💼 Happy selling and growing your business!"
    p.font.size = Pt(26)
    p.font.color.rgb = COLORS['white']
    p.alignment = PP_ALIGN.CENTER
    p.line_spacing = 1.4

    add_home_button(slide)

    # Save presentation
    output_file = "Shop_Management_System_User_Guide_EN.pptx"
    prs.save(output_file)

    try:
        print(f"\nPresentation created successfully: {output_file}")
    except UnicodeEncodeError:
        print("\nPresentation created successfully!")

    print(f"Total slides: {len(prs.slides)}")

    return output_file


if __name__ == "__main__":
    print("=" * 70)
    print("Enhanced Shop Management System - English User Guide")
    print("With Colors, Emojis, and Navigation Buttons")
    print("=" * 70)
    print()

    output_file = generate_presentation()

    print()
    print("=" * 70)
    print("Next steps:")
    print("1. Open the PPTX file in PowerPoint or Google Slides")
    print("2. Click the Home button on any slide to jump to Table of Contents")
    print("3. Add screenshots where placeholders are indicated")
    print("4. Customize colors if needed")
    print("5. Add your company logo")
    print("=" * 70)
