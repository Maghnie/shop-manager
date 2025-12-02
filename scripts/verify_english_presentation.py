#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""
Verify the English presentation
"""

from pptx import Presentation


def verify_presentation(filename):
    """Verify the presentation and show details"""
    print("=" * 70)
    print("English Presentation Verification")
    print("=" * 70)
    print()

    try:
        prs = Presentation(filename)
        print("File loaded successfully!")
        print()
        print(f"Total slides: {len(prs.slides)}")
        print(f"Slide dimensions: {prs.slide_width.inches:.1f}\" x {prs.slide_height.inches:.1f}\"")
        print()
        print("Slide Overview:")
        print("-" * 70)

        section_count = 0
        content_count = 0

        for i, slide in enumerate(prs.slides, 1):
            title = ""
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    text = shape.text
                    # Get first line of text
                    first_line = text.split('\n')[0][:50]
                    if first_line and not first_line.startswith("Screenshot"):
                        title = first_line
                        break

            slide_type = ""
            if i == 1:
                slide_type = "[TITLE]"
            elif i == 2:
                slide_type = "[INTRO]"
            elif i == 3:
                slide_type = "[TOC]"
            elif "Section" in title:
                slide_type = "[SECTION DIVIDER]"
                section_count += 1
            elif i == len(prs.slides):
                slide_type = "[CLOSING]"
            else:
                slide_type = "[CONTENT]"
                content_count += 1

            print(f"{i:2d}. {slide_type:20s} {title}")

        print()
        print("-" * 70)
        print(f"Summary:")
        print(f"  - Total slides: {len(prs.slides)}")
        print(f"  - Section dividers: {section_count}")
        print(f"  - Content slides: {content_count}")
        print(f"  - Special slides: 3 (Title, Intro, TOC) + 1 (Closing)")
        print()
        print("=" * 70)
        print("Verification complete! Presentation is ready.")
        print("=" * 70)

    except Exception as e:
        print(f"Error: {e}")


if __name__ == "__main__":
    verify_presentation("Shop_Management_System_User_Guide_EN.pptx")
