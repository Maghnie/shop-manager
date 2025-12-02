#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""
Verify the generated PowerPoint presentation
"""

from pptx import Presentation


def verify_presentation(filename):
    """Verify the presentation and show details"""
    print("=" * 60)
    print("Presentation Verification")
    print("=" * 60)
    print()

    try:
        prs = Presentation(filename)
        print("File loaded successfully!")
        print()
        print(f"Total slides: {len(prs.slides)}")
        print(f"Slide dimensions: {prs.slide_width} x {prs.slide_height}")
        print()
        print("Slide Titles:")
        print("-" * 60)

        for i, slide in enumerate(prs.slides, 1):
            title = "No title"
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text:
                    title = shape.text[:60]
                    break

            try:
                print(f"{i:2d}. {title}")
            except UnicodeEncodeError:
                print(f"{i:2d}. [Arabic title - encoding issue in terminal]")

        print()
        print("=" * 60)
        print("Verification complete!")
        print("=" * 60)

    except Exception as e:
        print(f"Error: {e}")


if __name__ == "__main__":
    verify_presentation("دليل_المستخدم_نظام_إدارة_المحل.pptx")
