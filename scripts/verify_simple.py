#!/usr/bin/env python
# -*- coding: utf-8 -*-
"""
Simple verification
"""

from pptx import Presentation


def verify_presentation(filename):
    """Verify the presentation"""
    print("=" * 70)
    print("Presentation Verification")
    print("=" * 70)
    print()

    try:
        prs = Presentation(filename)
        print(f"File: {filename}")
        print(f"Status: Loaded successfully!")
        print(f"Total slides: {len(prs.slides)}")
        print(f"Dimensions: {prs.slide_width.inches:.1f}\" x {prs.slide_height.inches:.1f}\"")
        print()

        # Count shapes with home buttons
        home_button_count = 0
        for slide in prs.slides:
            for shape in slide.shapes:
                if hasattr(shape, "text_frame"):
                    try:
                        if "Home" in shape.text_frame.text:
                            home_button_count += 1
                            break
                    except:
                        pass

        print(f"Slides with Home button: {home_button_count}")
        print()
        print("=" * 70)
        print("Presentation is ready for use!")
        print("=" * 70)

    except Exception as e:
        print(f"Error: {e}")


if __name__ == "__main__":
    import sys
    if len(sys.argv) > 1:
        verify_presentation(sys.argv[1])
    else:
        verify_presentation("Shop_Management_System_User_Guide_EN.pptx")
